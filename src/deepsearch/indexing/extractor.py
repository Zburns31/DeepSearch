"""
Text extraction from various file types
"""

import magic
from typing import Optional

# Document processing imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from .logger import create_extraction_logger, IndexingLogger


class TextExtractor:
    """Handles text extraction from various file types"""

    def __init__(self, logger: Optional[IndexingLogger] = None):
        self.logger = logger or create_extraction_logger()
        self.magic_detector = magic.Magic(mime=True)
        self._check_dependencies()

    def _check_dependencies(self):
        """Check for optional dependencies and log warnings"""
        missing_deps = []

        if PyPDF2 is None:
            missing_deps.append("PyPDF2 (for PDF support)")
        if Document is None:
            missing_deps.append("python-docx (for Word document support)")
        if openpyxl is None:
            missing_deps.append("openpyxl (for Excel support)")
        if Presentation is None:
            missing_deps.append("python-pptx (for PowerPoint support)")

        if missing_deps:
            self.logger.warning(
                f"Missing optional dependencies: {', '.join(missing_deps)}. "
                "Some file types will not be processed."
            )

    def extract_text(self, file_path: str) -> tuple[str, str]:
        """
        Extract text content from file

        Returns:
            tuple: (extracted_text, mime_type)
        """
        try:
            mime_type = self.magic_detector.from_file(file_path)

            # Route to appropriate extractor based on MIME type
            if mime_type == "application/pdf":
                return self._extract_pdf(file_path), mime_type
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                return self._extract_docx(file_path), mime_type
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ):
                return self._extract_xlsx(file_path), mime_type
            elif (
                mime_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                return self._extract_pptx(file_path), mime_type
            elif mime_type.startswith("text/"):
                return self._extract_text_file(file_path), mime_type
            elif mime_type in ["application/javascript", "application/json"]:
                return self._extract_text_file(file_path), mime_type
            elif file_path.endswith(".ipynb"):
                return self._extract_jupyter_notebook(file_path), mime_type
            else:
                return "", mime_type

        except Exception as e:
            self.logger.log_file_failed(file_path, e)
            return "", "unknown"

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF files"""
        if PyPDF2 is None:
            raise ImportError("PyPDF2 not available for PDF extraction")

        try:
            with open(file_path, "rb") as file:
                reader = PyPDF2.PdfReader(file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()
        except Exception as e:
            self.logger.debug(f"PDF extraction failed for {file_path}: {e}")
            return ""

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from Word documents"""
        if Document is None:
            raise ImportError("python-docx not available for Word document extraction")

        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return text.strip()
        except Exception as e:
            self.logger.debug(f"DOCX extraction failed for {file_path}: {e}")
            return ""

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from Excel files"""
        if openpyxl is None:
            raise ImportError("openpyxl not available for Excel extraction")

        try:
            workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(cell) for cell in row if cell is not None])
                    if row_text.strip():
                        text += row_text + "\n"
            workbook.close()
            return text.strip()
        except Exception as e:
            self.logger.debug(f"XLSX extraction failed for {file_path}: {e}")
            return ""

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint presentations"""
        if Presentation is None:
            raise ImportError("python-pptx not available for PowerPoint extraction")

        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        shape_text = getattr(shape, "text", "")
                        if shape_text and shape_text.strip():
                            text += shape_text + "\n"
            return text.strip()
        except Exception as e:
            self.logger.debug(f"PPTX extraction failed for {file_path}: {e}")
            return ""

    def _extract_text_file(self, file_path: str) -> str:
        """Extract text from plain text files"""
        try:
            # Try different encodings
            encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]

            for encoding in encodings:
                try:
                    with open(file_path, "r", encoding=encoding) as file:
                        content = file.read()
                        # Limit size to prevent memory issues
                        if len(content) > 1024 * 1024:  # 1MB limit
                            content = content[: 1024 * 1024] + "... [truncated]"
                        return content
                except UnicodeDecodeError:
                    continue

            # If all encodings fail, try with error handling
            with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
                content = file.read()
                if len(content) > 1024 * 1024:
                    content = content[: 1024 * 1024] + "... [truncated]"
                return content

        except Exception as e:
            self.logger.debug(f"Text extraction failed for {file_path}: {e}")
            return ""

    def _extract_jupyter_notebook(self, file_path: str) -> str:
        """Extract text from Jupyter notebook files"""
        try:
            import json

            with open(file_path, "r", encoding="utf-8") as file:
                notebook = json.load(file)

            text = ""

            # Extract text from cells
            for cell in notebook.get("cells", []):
                cell_type = cell.get("cell_type", "")

                if cell_type == "markdown":
                    # Extract markdown content
                    source = cell.get("source", [])
                    if isinstance(source, list):
                        text += "".join(source) + "\n\n"
                    else:
                        text += str(source) + "\n\n"

                elif cell_type == "code":
                    # Extract code content
                    source = cell.get("source", [])
                    if isinstance(source, list):
                        text += "".join(source) + "\n\n"
                    else:
                        text += str(source) + "\n\n"

                    # Extract output text if available
                    outputs = cell.get("outputs", [])
                    for output in outputs:
                        if output.get("output_type") == "stream":
                            text_output = output.get("text", [])
                            if isinstance(text_output, list):
                                text += "".join(text_output) + "\n"
                            else:
                                text += str(text_output) + "\n"

            return text.strip()

        except Exception as e:
            self.logger.debug(
                f"Jupyter notebook extraction failed for {file_path}: {e}"
            )
            return ""

    def can_extract(self, mime_type: str) -> bool:
        """Check if we can extract text from this MIME type"""
        extractable_types = {
            "application/pdf",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/javascript",
            "application/json",
        }

        return mime_type in extractable_types or mime_type.startswith("text/")
