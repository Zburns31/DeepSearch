"""
Standalone text extraction functions for multiprocessing
"""

import magic
from pathlib import Path

# Document processing imports
try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def extract_text_standalone(file_path: str) -> tuple[str, str]:
    """
    Standalone text extraction function that can be used with ProcessPoolExecutor

    Returns:
        tuple: (extracted_text, mime_type)
    """
    try:
        # Create a fresh magic detector for each call
        magic_detector = magic.Magic(mime=True)
        mime_type = magic_detector.from_file(file_path)

        # Route to appropriate extractor based on MIME type
        if mime_type == "application/pdf":
            return _extract_pdf_standalone(file_path), mime_type
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            return _extract_docx_standalone(file_path), mime_type
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        ):
            return _extract_xlsx_standalone(file_path), mime_type
        elif (
            mime_type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            return _extract_pptx_standalone(file_path), mime_type
        elif mime_type.startswith("text/") or mime_type in [
            "application/javascript",
            "application/json",
        ]:
            return _extract_text_file_standalone(file_path), mime_type
        else:
            return "", mime_type

    except Exception as e:
        print(f"Standalone extraction failed for {file_path}: {e}")
        return "", "unknown"


def _extract_pdf_standalone(file_path: str) -> str:
    """Extract text from PDF files"""
    if PyPDF2 is None:
        return ""

    try:
        with open(file_path, "rb") as file:
            reader = PyPDF2.PdfReader(file)
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            return text.strip()
    except Exception:
        return ""


def _extract_docx_standalone(file_path: str) -> str:
    """Extract text from Word documents"""
    if Document is None:
        return ""

    try:
        doc = Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip()
    except Exception:
        return ""


def _extract_xlsx_standalone(file_path: str) -> str:
    """Extract text from Excel files"""
    if openpyxl is None:
        return ""

    try:
        workbook = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text = ""
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join([str(cell) for cell in row if cell is not None])
                if row_text.strip():
                    text += row_text + "\n"
        workbook.close()
        return text.strip()
    except Exception:
        return ""


def _extract_pptx_standalone(file_path: str) -> str:
    """Extract text from PowerPoint presentations"""
    if Presentation is None:
        return ""

    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = getattr(shape, "text", "")
                    if shape_text and shape_text.strip():
                        text += shape_text + "\n"
        return text.strip()
    except Exception:
        return ""


def _extract_text_file_standalone(file_path: str) -> str:
    """Extract text from plain text files"""
    try:
        # Try different encodings
        encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    content = file.read()
                    # Limit size to prevent memory issues
                    if len(content) > 1024 * 1024:  # 1MB limit
                        content = content[: 1024 * 1024] + "... [truncated]"
                    return content
            except UnicodeDecodeError:
                continue

        # If all encodings fail, try with error handling
        with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
            content = file.read()
            if len(content) > 1024 * 1024:
                content = content[: 1024 * 1024] + "... [truncated]"
            return content

    except Exception:
        return ""
